import os
import tempfile
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

def handle_file_conversion(file, target_format):
    import shutil
    import pythoncom
    import win32com.client
    
    ext = os.path.splitext(file.filename)[1].lower()
    td = tempfile.gettempdir()
    ip = os.path.join(td, file.filename)
    file.save(ip)
    of = os.path.splitext(file.filename)[0] + target_format
    op = os.path.join(td, of)
    
    pythoncom.CoInitialize()
    try:
        converted = False
        if ext == '.csv' and target_format == '.xlsx':
            pd.read_csv(ip).to_excel(op, index=False)
            converted = True
        elif ext == '.xlsx' and target_format == '.csv':
            pd.read_excel(ip).to_csv(op, index=False)
            converted = True
        elif ext in ['.xlsx', '.xls'] and target_format == '.pdf':
            excel = win32com.client.Dispatch("Excel.Application")
            excel.DisplayAlerts = False
            wb = excel.Workbooks.Open(ip)
            wb.ExportAsFixedFormat(0, op)
            wb.Close()
            excel.Quit()
            converted = True
        elif ext in ['.docx', '.doc'] and target_format == '.pdf':
            word = win32com.client.Dispatch("Word.Application")
            word.DisplayAlerts = False
            doc = word.Documents.Open(ip)
            doc.SaveAs(op, 17)
            doc.Close()
            word.Quit()
            converted = True
        elif ext in ['.pptx', '.ppt'] and target_format == '.pdf':
            ppt = win32com.client.Dispatch("PowerPoint.Application")
            ppt.DisplayAlerts = False
            pres = ppt.Presentations.Open(ip, WithWindow=False)
            pres.SaveAs(op, 32)
            pres.Close()
            ppt.Quit()
            converted = True
        elif ext == '.xlsx' and target_format == '.pptx':
            ex = pd.read_excel(ip, sheet_name=None)
            prs = Presentation()
            for sname, df in ex.items():
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = sname
                tb = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(5))
                tb.text_frame.text = df.head(50).to_string(index=False)[:1000]
            prs.save(op)
            converted = True
        elif ext == '.pptx' and target_format == '.xlsx':
            prs = Presentation(ip)
            rows = [{"Slide": i+1, "Text": "\n".join([s.text for s in cl.shapes if hasattr(s,"text")])} for i, cl in enumerate(prs.slides)]
            pd.DataFrame(rows).to_excel(op, index=False)
            converted = True
            
        if not converted:
            raise Exception("지원되지 않는 변환이거나 호환되지 않는 포맷입니다. (단순 확장자 변경으로 인한 손상 방지)")

        return {"message": "✅ 완료", "path": op, "filename": of}
    except Exception as e: 
        raise Exception(str(e))
    finally:
        pythoncom.CoUninitialize()
        try:
            if os.path.exists(ip) and ip != op:
                os.remove(ip)
        except:
            pass

def batch_rename_extension(dir_path, old_ext, new_ext):
    dir_path = str(dir_path).strip('\"\'').strip()
    if not os.path.isdir(dir_path): raise Exception("존재하지 않는 폴더 경로입니다.")
    old_ext = ('.' + old_ext.strip('.'))
    new_ext = ('.' + new_ext.strip('.'))
    count = 0
    for fn in os.listdir(dir_path):
        if fn.lower().endswith(old_ext.lower()):
            new_fn = os.path.splitext(fn)[0] + new_ext
            os.rename(os.path.join(dir_path, fn), os.path.join(dir_path, new_fn))
            count += 1
    return count

def lock_office_file(file, password):
    import pythoncom
    import win32com.client
    td = tempfile.gettempdir()
    ip = os.path.join(td, "temp_" + file.filename)
    file.save(ip)
    
    of = "locked_" + file.filename
    op = os.path.join(td, of)
    
    ext = os.path.splitext(file.filename)[1].lower()
    
    pythoncom.CoInitialize()
    try:
        if ext in ['.xlsx', '.xls']:
            excel = win32com.client.Dispatch("Excel.Application")
            excel.DisplayAlerts = False
            wb = excel.Workbooks.Open(ip)
            wb.Password = password
            fmt = 51 if ext == '.xlsx' else 56
            wb.SaveAs(op, fmt)
            wb.Close()
            excel.Quit()
        elif ext in ['.docx', '.doc']:
            word = win32com.client.Dispatch("Word.Application")
            word.DisplayAlerts = False
            doc = word.Documents.Open(ip)
            doc.Password = password
            fmt = 12 if ext == '.docx' else 0
            doc.SaveAs(op, fmt)
            doc.Close()
            word.Quit()
        elif ext in ['.pptx', '.ppt']:
            ppt = win32com.client.Dispatch("PowerPoint.Application")
            ppt.DisplayAlerts = False
            pres = ppt.Presentations.Open(ip, WithWindow=False)
            pres.Password = password
            pres.SaveAs(op)
            pres.Close()
            ppt.Quit()
        else:
            raise Exception("지원되지 않는 오피스 파일 형식입니다. (xlsx, docx, pptx 등)")
    except Exception as e:
        raise Exception(f"암호 설정 실패: {str(e)}")
    finally:
        pythoncom.CoUninitialize()
        try:
            if os.path.exists(ip):
                os.remove(ip)
        except:
            pass
            
    if not os.path.exists(op):
        raise Exception("암호화된 파일이 정상적으로 생성되지 않았습니다.")
        
    return {"message": "🔒 보안 암호 설정 완료", "path": op}
